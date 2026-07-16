"""OfficeExtractor — extrae metadatos de documentos Office (DOCX, XLSX, PPTX).

Dependencias:
  - python-docx (para DOCX)
  - openpyxl (para XLSX)
  - pptx (para PPTX)

Todas son opcionales con degradación graceful individual.
Si una librería no está instalada, el extractor omite ese formato.
"""

from __future__ import annotations

import logging
import time
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from knowledge.engine.extractors.base import (
    ExtractionResult,
    _check_import,
    _hash_stream,
    get_registry,
)
from knowledge.engine.ontology.internal import AssetSource, AssetType, KnowledgeAsset

log = logging.getLogger("ura.knowledge.extractors.office")

_OFFICE_MIMES = [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
]

MAX_OFFICE_SIZE = 200 * 1024 * 1024

_HAS_DOCX = _check_import("docx", "python-docx")
_HAS_OPENPYXL = _check_import("openpyxl", "openpyxl")
_HAS_PPTX = _check_import("pptx", "python-pptx")


class OfficeExtractor:
    """Extractor para documentos Office (DOCX, XLSX, PPTX).

    Uso:
        extractor = OfficeExtractor()
        result = extractor.extract(source)
    """

    id: str = "office"
    version: str = "1.0.0"
    supported_mime_types: list[str] = _OFFICE_MIMES
    cost: str = "O(n)"

    def extract(self, source: AssetSource) -> ExtractionResult:
        t0 = time.monotonic()
        path_str = source.location
        path = Path(path_str)

        if not path.exists():
            return ExtractionResult(
                errors=[f"File not found: {path_str}"],
                duration_ms=(time.monotonic() - t0) * 1000,
            )

        try:
            content_sha256, size = _hash_stream(path_str)

            if size > MAX_OFFICE_SIZE:
                return ExtractionResult(
                    errors=[f"File too large: {size} bytes (max {MAX_OFFICE_SIZE})"],
                    duration_ms=(time.monotonic() - t0) * 1000,
                )

            now = datetime.now(UTC).isoformat()
            metadata: dict[str, Any] = {
                "size": size,
                "content_sha256": content_sha256,
                "format": path.suffix.lower().lstrip("."),
                "_extractor": self.id,
                "_extractor_version": self.version,
                "wraps": f"source:{path_str}",
                "extracted_at": now,
            }

            asset_type, warnings = self._extract_by_extension(path, metadata)
            if metadata.get("_degraded_reason"):
                log.warning("Degraded extraction for %s: %s", path_str, metadata["_degraded_reason"])

            asset = KnowledgeAsset(
                asset_id=content_sha256[:16],
                asset_type=asset_type,
                metadata=metadata,
                source=source,
                quality=_compute_office_quality(metadata),
                created_at=now,
                updated_at=now,
            )

            return ExtractionResult(
                asset=asset,
                warnings=warnings,
                duration_ms=(time.monotonic() - t0) * 1000,
            )

        except Exception as exc:
            log.exception("Office extraction error for %s", path_str)
            return ExtractionResult(
                errors=[f"Extraction error: {exc}"],
                duration_ms=(time.monotonic() - t0) * 1000,
            )

    def _extract_by_extension(self, path: Path, metadata: dict[str, Any]) -> tuple[AssetType, list[str]]:
        suffix = path.suffix.lower()

        if suffix == ".docx":
            atype, warns = self._extract_docx(path, metadata)
            return atype, warns
        if suffix == ".xlsx":
            atype, warns = self._extract_xlsx(path, metadata)
            return atype, warns
        if suffix == ".pptx":
            atype, warns = self._extract_pptx(path, metadata)
            return atype, warns

        metadata["_degraded"] = True
        metadata["_degraded_reason"] = f"Unsupported extension: {suffix}"
        return AssetType.OFFICE_DOC, []

    def _extract_docx(self, path: Path, metadata: dict[str, Any]) -> tuple[AssetType, list[str]]:
        if not _HAS_DOCX:
            metadata["_degraded"] = True
            metadata["_degraded_reason"] = "python-docx not installed"
            return AssetType.OFFICE_DOC, []

        import docx

        doc = docx.Document(str(path))

        cp = doc.core_properties
        for key, attr in (
            ("title", "title"),
            ("author", "author"),
            ("subject", "subject"),
            ("category", "category"),
            ("comments", "comments"),
            ("keywords", "keywords"),
            ("last_modified_by", "last_modified_by"),
            ("created", "created_at"),
            ("modified", "modified_at"),
        ):
            val = getattr(cp, attr, None)
            if val:
                val_str = str(val)
                if val_str:
                    metadata[f"office_{key}"] = val_str

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        metadata["paragraph_count"] = len(paragraphs)
        metadata["word_count"] = sum(len(p.split()) for p in paragraphs)
        metadata["text_preview"] = paragraphs[0][:300] if paragraphs else ""

        tables_count = len(doc.tables)
        metadata["tables_count"] = tables_count
        if tables_count > 0:
            total_rows = sum(len(t.rows) for t in doc.tables)
            total_cells = sum(len(t.rows) * len(t.columns) for t in doc.tables)
            metadata["tables_rows_total"] = total_rows
            metadata["tables_cells_total"] = total_cells

        sections_count = len(doc.sections)
        metadata["sections_count"] = sections_count

        log.info("Extracted DOCX: %s (%d paragraphs, %d tables)", path, len(paragraphs), tables_count)
        return AssetType.OFFICE_DOC, []

    def _extract_xlsx(self, path: Path, metadata: dict[str, Any]) -> tuple[AssetType, list[str]]:
        if not _HAS_OPENPYXL:
            metadata["_degraded"] = True
            metadata["_degraded_reason"] = "openpyxl not installed"
            return AssetType.OFFICE_SHEET, []

        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        try:
            props = wb.properties
            if props:
                for key in ("title", "subject", "keywords", "category", "description", "creator"):
                    val = getattr(props, key, None)
                    if val and str(val).strip():
                        metadata[f"office_{key}"] = str(val)

            metadata["sheet_names"] = wb.sheetnames
            metadata["sheet_count"] = len(wb.sheetnames)

            total_rows = 0
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                total_rows += ws.max_row or 0

            metadata["rows_total"] = total_rows
            warns = ["Row counts may be approximate (read_only mode)"]

            log.info("Extracted XLSX: %s (%d sheets, ~%d rows)", path, len(wb.sheetnames), total_rows)
        finally:
            wb.close()

        return AssetType.OFFICE_SHEET, warns

    def _extract_pptx(self, path: Path, metadata: dict[str, Any]) -> tuple[AssetType, list[str]]:
        if not _HAS_PPTX:
            metadata["_degraded"] = True
            metadata["_degraded_reason"] = "python-pptx not installed"
            return AssetType.OFFICE_SLIDE, []

        from pptx import Presentation

        prs = Presentation(str(path))

        cp = prs.core_properties
        for key in (
            "title",
            "author",
            "subject",
            "keywords",
            "comments",
            "category",
            "last_modified_by",
            "created",
            "modified",
        ):
            val = getattr(cp, key, None)
            if val and str(val).strip():
                metadata[f"office_{key}"] = str(val)

        metadata["slide_count"] = len(prs.slides)
        metadata["slide_width"] = prs.slide_width
        metadata["slide_height"] = prs.slide_height

        shapes_total = 0
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                shapes_total += 1
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

        metadata["shapes_total"] = shapes_total
        if texts:
            preview = texts[0][:300]
            metadata["text_preview"] = preview

        log.info("Extracted PPTX: %s (%d slides, %d shapes)", path, len(prs.slides), shapes_total)
        return AssetType.OFFICE_SLIDE, []


def _compute_office_quality(metadata: dict[str, Any]) -> float:
    q = 0.3
    has_content = (
        metadata.get("paragraph_count", 0) > 0
        or metadata.get("sheet_count", 0) > 0
        or metadata.get("slide_count", 0) > 0
    )
    if has_content:
        q += 0.15
    if metadata.get("office_title"):
        q += 0.15
    if metadata.get("office_author"):
        q += 0.1
    if metadata.get("word_count", 0) > 50:
        q += 0.15
    if metadata.get("tables_count", 0) > 0:
        q += 0.15
    if not metadata.get("_degraded", False):
        q += 0.1
    return min(q, 1.0)


_registry = get_registry()
_registry.register(OfficeExtractor())
