"""Tests for extractors (Fase 5)."""

import hashlib
import shutil
import subprocess
from pathlib import Path

import pytest

from knowledge.engine.extractors.base import (
    _check_import,
    _hash_stream,
    get_registry,
)
from knowledge.engine.ontology.internal import AssetSource, AssetType

# ── Fixtures ─────────────────────────────────────────────────────────────────


@pytest.fixture
def pdf_sample(tmp_path: Path) -> Path:
    """Crea un PDF de prueba con una página de texto."""
    p = tmp_path / "sample.pdf"
    fitz = pytest.importorskip("fitz")
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF World", fontname="helv", fontsize=12)
    doc.save(str(p))
    doc.close()
    return p


@pytest.fixture
def pdf_sample_2pages(tmp_path: Path) -> Path:
    """Crea un PDF de prueba con 2 páginas."""
    p = tmp_path / "sample2.pdf"
    fitz = pytest.importorskip("fitz")
    doc = fitz.open()
    doc.new_page()
    doc.new_page()
    doc.save(str(p))
    doc.close()
    return p


# ── Helpers tests ────────────────────────────────────────────────────────────


class TestHelpers:
    def test_hash_stream(self, tmp_path: Path):
        p = tmp_path / "test.bin"
        p.write_bytes(b"hello world")
        h, size = _hash_stream(p)
        expected = hashlib.sha256(b"hello world").hexdigest()
        assert h == expected
        assert size == 11

    def test_hash_stream_empty(self):
        h, size = _hash_stream("/dev/null")
        assert h == hashlib.sha256(b"").hexdigest()
        assert size == 0

    def test_hash_stream_not_found(self, tmp_path: Path):
        with pytest.raises(FileNotFoundError):
            _hash_stream(str(tmp_path / "nonexistent_file_xyz"))

    def test_check_import_found(self):
        assert _check_import("os") is True

    def test_check_import_not_found(self):
        assert _check_import("nonexistent_module_xyz") is False


# ── PdfExtractor tests ──────────────────────────────────────────────────────


class TestPdfExtractor:
    def test_extract_basic(self, pdf_sample: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type == AssetType.PDF
        assert result.asset.metadata.get("pages") == 1
        assert result.asset.metadata.get("has_text") is True
        assert result.asset.metadata["_extractor"] == "pdf"
        assert result.duration_ms > 0

    def test_metadata_fields(self, pdf_sample: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        meta = result.asset.metadata
        for required in [
            "pages",
            "content_sha256",
            "size",
            "has_text",
            "text_length",
            "text_preview",
            "_extractor",
            "_extractor_version",
            "extracted_at",
            "wraps",
        ]:
            assert required in meta, f"Missing metadata field: {required}"

    def test_determinism(self, pdf_sample: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        r1 = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        r2 = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        assert r1.asset.asset_id == r2.asset.asset_id

    def test_asset_id_is_content_hash(self, pdf_sample: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        content_hash = result.asset.metadata["content_sha256"]
        assert result.asset.asset_id == content_hash[:16]

    def test_quality_computed(self, pdf_sample: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(pdf_sample)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_file_not_found(self, tmp_path: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "nonexistent.pdf")))
        assert len(result.errors) > 0
        assert "not found" in result.errors[0].lower()
        assert result.asset is None

    def test_multipage(self, pdf_sample_2pages: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(pdf_sample_2pages)))
        assert result.asset.metadata["pages"] == 2

    def test_no_exception_on_corrupt(self, tmp_path: Path):
        from knowledge.engine.extractors.pdf import PdfExtractor

        p = tmp_path / "corrupt.pdf"
        p.write_bytes(b"not a real pdf file content here")
        ext = PdfExtractor()
        result = ext.extract(AssetSource("filesystem", str(p)))
        assert result.errors or result.asset is not None
        if result.asset:
            assert result.asset.asset_type == AssetType.PDF

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("pdf")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert "application/pdf" in ext.supported_mime_types


# ── ImageExtractor tests ──────────────────────────────────────────────────────


@pytest.fixture
def image_sample_jpg(tmp_path: Path) -> Path:
    """Crea una imagen JPEG de prueba."""
    pytest.importorskip("PIL")
    from PIL import Image

    img = Image.new("RGB", (100, 100), color="red")
    p = tmp_path / "sample.jpg"
    img.save(str(p), "JPEG")
    return p


@pytest.fixture
def image_sample_png(tmp_path: Path) -> Path:
    """Crea una imagen PNG de prueba."""
    pytest.importorskip("PIL")
    from PIL import Image

    img = Image.new("RGBA", (50, 50), color="blue")
    p = tmp_path / "sample.png"
    img.save(str(p), "PNG")
    return p


class TestImageExtractor:
    def test_extract_jpeg(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type == AssetType.IMAGE
        assert result.asset.metadata.get("width") == 100
        assert result.asset.metadata.get("height") == 100
        assert result.duration_ms > 0

    def test_extract_png(self, image_sample_png: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_png)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.metadata.get("width") == 50
        assert result.asset.metadata.get("height") == 50

    def test_metadata_fields(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        meta = result.asset.metadata
        for required in [
            "width",
            "height",
            "format",
            "mode",
            "content_sha256",
            "size",
            "_extractor",
            "_extractor_version",
            "extracted_at",
            "wraps",
        ]:
            assert required in meta, f"Missing metadata field: {required}"

    def test_determinism(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        r1 = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        r2 = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        assert r1.asset.asset_id == r2.asset.asset_id

    def test_asset_id_is_content_hash(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        assert result.asset.asset_id == result.asset.metadata["content_sha256"][:16]

    def test_thumbnail_created(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        thumb = result.asset.metadata.get("thumbnail", "")
        assert thumb
        assert Path(thumb).exists()

    def test_file_not_found(self, tmp_path: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "nonexistent.jpg")))
        assert len(result.errors) > 0
        assert "not found" in result.errors[0].lower()
        assert result.asset is None

    def test_no_exception_on_corrupt(self, tmp_path: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        p = tmp_path / "corrupt.jpg"
        p.write_bytes(b"not a real image file content here")
        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(p)))
        assert result.asset is not None  # degraded asset with basic metadata
        assert result.asset.metadata.get("_degraded", False)

    def test_quality_computed(self, image_sample_jpg: Path):
        from knowledge.engine.extractors.image import ImageExtractor

        ext = ImageExtractor()
        result = ext.extract(AssetSource("filesystem", str(image_sample_jpg)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("image")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert "image/jpeg" in ext.supported_mime_types


# ── OfficeExtractor tests ────────────────────────────────────────────────────


@pytest.fixture
def office_sample_docx(tmp_path: Path) -> Path:
    """Crea un documento DOCX de prueba."""
    pytest.importorskip("docx")
    import docx

    doc = docx.Document()
    doc.add_paragraph("Hello from DOCX test")
    doc.add_paragraph("Second paragraph with more content")
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    p = tmp_path / "test.docx"
    doc.save(str(p))
    return p


@pytest.fixture
def office_sample_pptx(tmp_path: Path) -> Path:
    """Crea una presentación PPTX de prueba."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = "Test Slide Content"
    p = tmp_path / "test.pptx"
    prs.save(str(p))
    return p


class TestOfficeExtractor:
    def test_extract_docx(self, office_sample_docx: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(office_sample_docx)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type.name == "OFFICE_DOC"
        assert result.asset.metadata.get("paragraph_count", 0) >= 2
        assert result.duration_ms > 0

    def test_extract_pptx(self, office_sample_pptx: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(office_sample_pptx)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type.name == "OFFICE_SLIDE"
        assert result.asset.metadata.get("slide_count", 0) >= 1

    def test_metadata_fields_docx(self, office_sample_docx: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(office_sample_docx)))
        meta = result.asset.metadata
        for required in [
            "paragraph_count",
            "word_count",
            "format",
            "content_sha256",
            "size",
            "_extractor",
            "extracted_at",
        ]:
            assert required in meta, f"Missing: {required}"

    def test_determinism(self, office_sample_docx: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        r1 = ext.extract(AssetSource("filesystem", str(office_sample_docx)))
        r2 = ext.extract(AssetSource("filesystem", str(office_sample_docx)))
        assert r1.asset.asset_id == r2.asset.asset_id

    def test_degradation_no_docx(self, tmp_path: Path, monkeypatch):
        import knowledge.engine.extractors.office as office_mod

        monkeypatch.setattr(office_mod, "_HAS_DOCX", False)
        from knowledge.engine.extractors.office import OfficeExtractor

        p = tmp_path / "test.docx"
        p.write_text("fake docx content")
        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(p)))
        assert result.asset is not None
        assert result.asset.metadata.get("_degraded", False)
        assert result.errors == []

    def test_file_not_found(self, tmp_path: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "nonexistent.docx")))
        assert len(result.errors) > 0
        assert "not found" in result.errors[0].lower()

    def test_quality_computed(self, office_sample_docx: Path):
        from knowledge.engine.extractors.office import OfficeExtractor

        ext = OfficeExtractor()
        result = ext.extract(AssetSource("filesystem", str(office_sample_docx)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("office")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert len(ext.supported_mime_types) == 3


# ── AudioExtractor tests ─────────────────────────────────────────────────────


@pytest.fixture
def audio_sample_wav(tmp_path: Path) -> Path:
    """Crea un archivo WAV de prueba con tono simple."""
    import math
    import struct
    import wave

    sample_rate = 22050
    duration = 0.5
    freq = 440.0
    num_samples = int(sample_rate * duration)
    samples = []
    for i in range(num_samples):
        val = int(32767 * 0.3 * math.sin(2.0 * math.pi * freq * i / sample_rate))
        samples.append(val)

    p = tmp_path / "test.wav"
    with wave.open(str(p), "w") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sample_rate)
        wf.writeframes(struct.pack(f"<{len(samples)}h", *samples))
    return p


class TestAudioExtractor:
    def test_extract_wav(self, audio_sample_wav: Path):
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        result = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type == AssetType.AUDIO
        assert result.duration_ms > 0

    def test_metadata_fields(self, audio_sample_wav: Path):
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        result = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        meta = result.asset.metadata
        for required in [
            "format",
            "size",
            "content_sha256",
            "_extractor",
            "_extractor_version",
            "extracted_at",
            "wraps",
        ]:
            assert required in meta, f"Missing: {required}"

    def test_ffprobe_metadata(self, audio_sample_wav: Path):
        import shutil

        if not shutil.which("ffprobe"):
            pytest.skip("ffprobe binary not found")
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        result = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        meta = result.asset.metadata
        assert not meta.get("_degraded_ffprobe", False), "ffprobe was degraded"
        assert "audio_duration_sec" in meta, "Missing duration"
        assert "audio_codec" in meta, "Missing codec"

    def test_determinism(self, audio_sample_wav: Path):
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        r1 = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        r2 = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        assert r1.asset.asset_id == r2.asset.asset_id

    def test_file_not_found(self, tmp_path: Path):
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "nonexistent.wav")))
        assert len(result.errors) > 0
        assert "not found" in result.errors[0].lower()

    def test_quality_computed(self, audio_sample_wav: Path):
        from knowledge.engine.extractors.audio import AudioExtractor

        ext = AudioExtractor()
        result = ext.extract(AssetSource("filesystem", str(audio_sample_wav)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("audio")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert "audio/mp3" in ext.supported_mime_types


# ── VideoExtractor tests ─────────────────────────────────────────────────────


@pytest.fixture
def video_sample_mov(tmp_path: Path) -> Path:
    """Crea un vídeo MOV de prueba con ffmpeg."""
    if not shutil.which("ffmpeg"):
        pytest.skip("ffmpeg not available")
    p = tmp_path / "test.mov"
    cmd = [
        "ffmpeg",
        "-y",
        "-f",
        "lavfi",
        "-i",
        "testsrc=duration=0.5:size=64x64:rate=10",
        "-f",
        "mov",
        str(p),
    ]
    subprocess.run(cmd, capture_output=True, timeout=30, check=False)
    if not p.exists():
        pytest.skip("ffmpeg could not create test video")
    return p


class TestVideoExtractor:
    def test_extract_mov(self, video_sample_mov: Path):
        from knowledge.engine.extractors.video import VideoExtractor

        ext = VideoExtractor()
        result = ext.extract(AssetSource("filesystem", str(video_sample_mov)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type == AssetType.VIDEO
        assert result.duration_ms > 0

    def test_metadata_fields(self, video_sample_mov: Path):
        from knowledge.engine.extractors.video import VideoExtractor

        ext = VideoExtractor()
        result = ext.extract(AssetSource("filesystem", str(video_sample_mov)))
        meta = result.asset.metadata
        for required in [
            "format",
            "size",
            "content_sha256",
            "_extractor",
            "_extractor_version",
            "extracted_at",
            "wraps",
        ]:
            assert required in meta, f"Missing: {required}"

    def test_determinism(self, video_sample_mov: Path):
        from knowledge.engine.extractors.video import VideoExtractor

        ext = VideoExtractor()
        r1 = ext.extract(AssetSource("filesystem", str(video_sample_mov)))
        r2 = ext.extract(AssetSource("filesystem", str(video_sample_mov)))
        assert r1.asset.asset_id == r2.asset.asset_id

    def test_file_not_found(self, tmp_path: Path):
        from knowledge.engine.extractors.video import VideoExtractor

        ext = VideoExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "nonexistent.mp4")))
        assert len(result.errors) > 0
        assert "not found" in result.errors[0].lower()

    def test_quality_computed(self, video_sample_mov: Path):
        from knowledge.engine.extractors.video import VideoExtractor

        ext = VideoExtractor()
        result = ext.extract(AssetSource("filesystem", str(video_sample_mov)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("video")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert "video/mp4" in ext.supported_mime_types


# ── WebExtractor tests ────────────────────────────────────────────────────────


class TestWebExtractor:
    def test_ssrf_block_file_scheme(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "file:///etc/passwd"))
        assert len(result.errors) > 0
        assert "scheme" in result.errors[0].lower()

    def test_ssrf_block_localhost(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "http://localhost:8080/admin"))
        assert len(result.errors) > 0
        assert "blocked" in result.errors[0].lower()

    def test_ssrf_block_private_ip(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "http://10.0.0.1/config"))
        assert len(result.errors) > 0
        assert "blocked" in result.errors[0].lower()

    def test_ssrf_block_metadata(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "http://169.254.169.254/latest/meta-data/"))
        assert len(result.errors) > 0
        assert "cloud metadata" in result.errors[0].lower() or "blocked" in result.errors[0].lower()

    def test_ssrf_block_loopback(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "http://127.0.0.1:11434/"))
        assert len(result.errors) > 0
        assert "blocked" in result.errors[0].lower()

    def test_empty_url(self):
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", ""))
        assert len(result.errors) > 0

    def test_degradation_no_httpx(self, monkeypatch):
        import knowledge.engine.extractors.web as web_mod

        monkeypatch.setattr(web_mod, "_HAS_HTTPX", False)
        from knowledge.engine.extractors.web import WebExtractor

        ext = WebExtractor()
        result = ext.extract(AssetSource("api", "http://example.com"))
        assert result.asset is not None
        assert result.asset.metadata.get("_degraded", False)

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("web")
        assert ext is not None
        assert ext.version == "1.0.0"
        assert "text/html" in ext.supported_mime_types


# ── GitExtractor tests ────────────────────────────────────────────────────────


class TestGitExtractor:
    def test_extract_local(self, tmp_path: Path):
        """Crea un repo git local y extrae metadatos."""
        repo = tmp_path / "myrepo"
        repo.mkdir()
        (repo / "README.md").write_text("# Test Repo")
        subprocess.run(["git", "init"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(
            ["git", "config", "user.email", "test@test.com"], cwd=str(repo), capture_output=True, check=False
        )
        subprocess.run(["git", "config", "user.name", "Test"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "add", "."], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "commit", "-m", "Initial commit"], cwd=str(repo), capture_output=True, check=False)

        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(repo)))
        assert result.errors == []
        assert result.asset is not None
        assert result.asset.asset_type == AssetType.GIT_REPO
        assert result.asset.metadata.get("commit_count", 0) >= 1
        assert result.duration_ms > 0

    def test_metadata_fields(self, tmp_path: Path):
        repo = tmp_path / "metarepo"
        repo.mkdir()
        (repo / "file.txt").write_text("content")
        subprocess.run(["git", "init"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.email", "t@t.com"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.name", "T"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "add", "."], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "commit", "-m", "feat: add file"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "tag", "v1.0.0"], cwd=str(repo), capture_output=True, check=False)

        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(repo)))
        meta = result.asset.metadata
        for required in [
            "commit_count",
            "tag_count",
            "current_branch",
            "_extractor",
            "_extractor_version",
            "extracted_at",
            "wraps",
        ]:
            assert required in meta, f"Missing: {required}"
        assert meta["tag_count"] >= 1
        assert meta["commit_count"] >= 1

    def test_not_a_repo(self, tmp_path: Path):
        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(tmp_path / "notarepo")))
        assert len(result.errors) > 0

    def test_empty_location(self):
        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("api", ""))
        assert len(result.errors) > 0

    def test_readme_detected(self, tmp_path: Path):
        repo = tmp_path / "readme_repo"
        repo.mkdir()
        (repo / "README.md").write_text("# Project Title\n\nDescription here.")
        subprocess.run(["git", "init"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.email", "t@t.com"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.name", "T"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "add", "."], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "commit", "-m", "init"], cwd=str(repo), capture_output=True, check=False)

        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(repo)))
        assert result.asset.metadata.get("readme_preview", "").startswith("# Project Title")

    def test_quality_computed(self, tmp_path: Path):
        repo = tmp_path / "qrepo"
        repo.mkdir()
        (repo / "f.txt").write_text("x")
        subprocess.run(["git", "init"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.email", "t@t.com"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.name", "T"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "add", "."], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "commit", "-m", "init"], cwd=str(repo), capture_output=True, check=False)

        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(repo)))
        assert 0.0 <= result.asset.quality <= 1.0

    def test_exceeds_max_size(self, tmp_path: Path, monkeypatch):
        import knowledge.engine.extractors.git as git_mod

        monkeypatch.setattr(git_mod, "MAX_CLONE_SIZE", 1)
        repo = tmp_path / "largerepo"
        repo.mkdir()
        (repo / "bigfile.bin").write_bytes(b"x" * 100)
        subprocess.run(["git", "init"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.email", "t@t.com"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "config", "user.name", "T"], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "add", "."], cwd=str(repo), capture_output=True, check=False)
        subprocess.run(["git", "commit", "-m", "big"], cwd=str(repo), capture_output=True, check=False)
        from knowledge.engine.extractors.git import GitExtractor

        ext = GitExtractor()
        result = ext.extract(AssetSource("filesystem", str(repo)))
        assert len(result.errors) > 0
        assert "too large" in result.errors[0].lower() or "max" in result.errors[0].lower()

    def test_registered(self):
        registry = get_registry()
        ext = registry.get("git")
        assert ext is not None
        assert ext.version == "1.0.0"
