from pathlib import Path

import docx
import pptx


def extraer_office(ruta: Path) -> dict:
    suf = ruta.suffix.lower()
    if suf == ".docx":
        return _extraer_docx(ruta)
    if suf in (".pptx", ".ppt"):
        return _extraer_pptx(ruta)
    return _extraer_generico(ruta)


def _extraer_docx(ruta: Path) -> dict:
    doc = docx.Document(str(ruta))
    parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
    tablas = []
    for t in doc.tables:
        for row in t.rows:
            celdas = [cell.text for cell in row.cells if cell.text.strip()]
            if celdas:
                tablas.append(" | ".join(celdas))

    props = doc.core_properties
    return {
        "tipo": "office",
        "metadatos": {
            "formato": "docx",
            "titulo": props.title or ruta.stem,
            "autor": props.author or "",
            "parrafos": len(parrafos),
            "tablas": len(doc.tables),
            "tamano_bytes": ruta.stat().st_size,
        },
        "texto_plano": "\n".join(parrafos[:500] + tablas[:50]),
        "ruta": str(ruta),
    }


def _extraer_pptx(ruta: Path) -> dict:
    prs = pptx.Presentation(str(ruta))
    diapositivas = []
    for i, slide in enumerate(prs.slides):
        textos = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        textos.append(p.text.strip())
        if textos:
            diapositivas.append(f"[Slide {i+1}] {' | '.join(textos[:5])}")
        if i >= 30:
            diapositivas.append(f"... ({len(prs.slides) - 31} slides restantes)")
            break

    props = prs.core_properties
    return {
        "tipo": "office",
        "metadatos": {
            "formato": "pptx",
            "titulo": props.title or ruta.stem,
            "autor": props.author or "",
            "slides": len(prs.slides),
            "tamano_bytes": ruta.stat().st_size,
        },
        "texto_plano": "\n".join(diapositivas),
        "ruta": str(ruta),
    }


def _extraer_generico(ruta: Path) -> dict:
    return {
        "tipo": "office",
        "metadatos": {
            "formato": ruta.suffix.lower().lstrip("."),
            "titulo": ruta.stem,
            "tamano_bytes": ruta.stat().st_size,
        },
        "texto_plano": ruta.read_text(errors="replace")[:10000],
        "ruta": str(ruta),
    }
